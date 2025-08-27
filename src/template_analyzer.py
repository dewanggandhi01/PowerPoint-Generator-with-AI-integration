"""
Template Analyzer module for extracting styles and layouts from PowerPoint templates.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import os
import logging
from typing import Dict, List, Any, Optional

logger = logging.getLogger(__name__)

class TemplateAnalyzer:
    """Analyzer for PowerPoint template files to extract styling and layout information."""
    
    def __init__(self, template_path: str):
        self.template_path = template_path
        self.presentation = None
        self.theme_info = {}
        self.layout_info = {}
        self.images = []
        self._load_template()
    
    def _load_template(self):
        """Load the PowerPoint template and analyze its structure."""
        try:
            self.presentation = Presentation(self.template_path)
            self._analyze_theme()
            self._analyze_layouts()
            self._extract_images()
            logger.info(f"Successfully loaded template: {self.template_path}")
        except Exception as e:
            logger.error(f"Error loading template: {str(e)}")
            raise
    
    def _analyze_theme(self):
        """Extract theme information from the template."""
        try:
            # Get the first slide to analyze colors and fonts
            if self.presentation.slides:
                slide = self.presentation.slides[0]
                
                # Extract color scheme (basic analysis)
                self.theme_info['colors'] = self._extract_colors(slide)
                self.theme_info['fonts'] = self._extract_fonts(slide)
            
            # Get slide dimensions
            self.theme_info['slide_width'] = self.presentation.slide_width
            self.theme_info['slide_height'] = self.presentation.slide_height
            
        except Exception as e:
            logger.warning(f"Error analyzing theme: {str(e)}")
            self._set_default_theme()
    
    def _analyze_layouts(self):
        """Analyze slide layouts available in the template."""
        try:
            layouts = []
            for i, layout in enumerate(self.presentation.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': layout.name,
                    'placeholders': []
                }
                
                # Analyze placeholders in this layout
                for placeholder in layout.placeholders:
                    placeholder_info = {
                        'idx': placeholder.placeholder_format.idx,
                        'type': placeholder.placeholder_format.type,
                        'left': placeholder.left,
                        'top': placeholder.top,
                        'width': placeholder.width,
                        'height': placeholder.height
                    }
                    layout_info['placeholders'].append(placeholder_info)
                
                layouts.append(layout_info)
            
            self.layout_info['layouts'] = layouts
            
        except Exception as e:
            logger.warning(f"Error analyzing layouts: {str(e)}")
            self.layout_info['layouts'] = []
    
    def _extract_images(self):
        """Extract images from the template for reuse."""
        try:
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        image_info = {
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height,
                            'image_data': shape.image.blob
                        }
                        self.images.append(image_info)
        
        except Exception as e:
            logger.warning(f"Error extracting images: {str(e)}")
    
    def _extract_colors(self, slide) -> Dict[str, str]:
        """Extract color scheme from a slide."""
        colors = {}
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                    if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        colors['primary'] = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                        break
        except Exception:
            pass
        
        # Set defaults if no colors found
        if 'primary' not in colors:
            colors['primary'] = '#1f4e79'
        colors['secondary'] = '#4472c4'
        colors['accent'] = '#70ad47'
        
        return colors
    
    def _extract_fonts(self, slide) -> Dict[str, str]:
        """Extract font information from a slide."""
        fonts = {'title': 'Calibri', 'body': 'Calibri'}
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                if 'title' not in fonts:
                                    fonts['title'] = run.font.name
                                fonts['body'] = run.font.name
                                break
        except Exception:
            pass
        
        return fonts
    
    def _set_default_theme(self):
        """Set default theme information."""
        self.theme_info = {
            'colors': {
                'primary': '#1f4e79',
                'secondary': '#4472c4',
                'accent': '#70ad47'
            },
            'fonts': {
                'title': 'Calibri',
                'body': 'Calibri'
            },
            'slide_width': Inches(10),
            'slide_height': Inches(7.5)
        }
    
    def get_best_layout_for_slide_type(self, slide_type: str) -> int:
        """
        Get the best layout index for a given slide type.
        
        Args:
            slide_type: Type of slide ('title', 'content', 'conclusion', 'section', 'comparison')
            
        Returns:
            Layout index to use
        """
        layouts = self.layout_info.get('layouts', [])
        
        if not layouts:
            return 0  # Default to first layout
        
        # Enhanced layout selection logic
        if slide_type == 'title':
            # Look for title slide layout (usually first)
            for layout in layouts:
                name_lower = layout['name'].lower()
                if any(keyword in name_lower for keyword in ['title', 'cover', 'intro']):
                    return layout['index']
            return 0
        
        elif slide_type == 'section':
            # Look for section header layouts
            for layout in layouts:
                name_lower = layout['name'].lower()
                if any(keyword in name_lower for keyword in ['section', 'divider', 'header', 'chapter']):
                    return layout['index']
            # Fallback to title-like layout
            return 0
        
        elif slide_type == 'content':
            # Look for content layouts with bullet points
            for layout in layouts:
                name_lower = layout['name'].lower()
                if any(keyword in name_lower for keyword in ['content', 'bullet', 'text', 'list']):
                    return layout['index']
            # Use layout 1 if available (usually content layout)
            return min(1, len(layouts) - 1)
        
        elif slide_type == 'comparison':
            # Look for two-column or comparison layouts
            for layout in layouts:
                name_lower = layout['name'].lower()
                if any(keyword in name_lower for keyword in ['two', 'comparison', 'column', 'vs']):
                    return layout['index']
            # Fallback to content layout
            return min(1, len(layouts) - 1)
        
        elif slide_type == 'conclusion':
            # Look for conclusion or thank you layouts
            for layout in layouts:
                name_lower = layout['name'].lower()
                if any(keyword in name_lower for keyword in ['conclusion', 'thank', 'end', 'summary']):
                    return layout['index']
            # Use title layout for conclusion
            return 0
        
        # Default fallback
        return min(1, len(layouts) - 1)
    
    def get_theme_colors(self) -> Dict[str, str]:
        """Get theme colors."""
        return self.theme_info.get('colors', {})
    
    def get_background_colors(self) -> Dict[str, str]:
        """Get background colors from the template for contrast calculation."""
        try:
            bg_colors = {}
            
            # Analyze first slide's background
            if self.presentation.slides:
                slide = self.presentation.slides[0]
                
                # Try to get background fill
                if hasattr(slide, 'background') and slide.background.fill:
                    fill = slide.background.fill
                    if hasattr(fill, 'fore_color') and fill.fore_color:
                        if hasattr(fill.fore_color, 'rgb'):
                            rgb = fill.fore_color.rgb
                            # Fix: Access RGB values correctly
                            try:
                                if hasattr(rgb, '__iter__') and len(rgb) >= 3:
                                    # RGB is a tuple or list
                                    bg_colors['primary'] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                else:
                                    # RGB is an RGBColor object
                                    bg_colors['primary'] = f"#{rgb:06x}"
                            except:
                                bg_colors['primary'] = '#FFFFFF'
                
                # Fallback: analyze slide master background
                if not bg_colors and hasattr(slide, 'slide_layout'):
                    layout = slide.slide_layout
                    if hasattr(layout, 'slide_master'):
                        master = layout.slide_master
                        if hasattr(master, 'background') and master.background.fill:
                            fill = master.background.fill
                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                if hasattr(fill.fore_color, 'rgb'):
                                    rgb = fill.fore_color.rgb
                                    try:
                                        if hasattr(rgb, '__iter__') and len(rgb) >= 3:
                                            bg_colors['primary'] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                        else:
                                            bg_colors['primary'] = f"#{rgb:06x}"
                                    except:
                                        bg_colors['primary'] = '#FFFFFF'
            
            # Default to white if no background found
            if not bg_colors:
                bg_colors['primary'] = '#FFFFFF'
                
            return bg_colors
            
        except Exception as e:
            logger.warning(f"Error getting background colors: {str(e)}")
            return {'primary': '#FFFFFF'}  # Default to white
    
    def get_theme_fonts(self) -> Dict[str, str]:
        """Get theme fonts."""
        return self.theme_info.get('fonts', {})
    
    def get_template_images(self) -> List[Dict]:
        """Get extracted images from template."""
        return self.images
    
    def get_slide_dimensions(self) -> tuple:
        """Get slide dimensions."""
        return (
            self.theme_info.get('slide_width', Inches(10)),
            self.theme_info.get('slide_height', Inches(7.5))
        )

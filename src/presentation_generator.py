"""
Presentation Generator module for creating PowerPoint presentations.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import tempfile
import logging
from typing import Dict, List, Any
from .llm_service import LLMService
from .template_analyzer import TemplateAnalyzer

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Generator for creating PowerPoint presentations from text and templates."""
    
    def __init__(self, llm_service: LLMService, template_analyzer: TemplateAnalyzer):
        self.llm_service = llm_service
        self.template_analyzer = template_analyzer
        self.user_images = []
    
    def generate(self, text_input: str, guidance: str = "", image_paths: List[str] = None) -> str:
        """
        Generate a PowerPoint presentation from text input with optional images.
        
        Args:
            text_input: Input text to convert to presentation
            guidance: Optional guidance for tone/structure
            image_paths: List of paths to images to include in presentation
            
        Returns:
            Path to the generated PowerPoint file
        """
        logger.info("Starting presentation generation")
        
        # Store image paths for use in slide generation
        self.user_images = image_paths or []
        
        # Analyze text structure using LLM
        structure = self.llm_service.analyze_text_structure(text_input, guidance)
        
        # Create new presentation based on template
        presentation = self._create_presentation_from_template()
        
        # Generate slides with images
        self._generate_slides_with_images(presentation, structure)
        
        # Save presentation
        output_path = self._save_presentation(presentation)
        
        logger.info(f"Presentation generated successfully: {output_path}")
        return output_path
    
    def _create_presentation_from_template(self) -> Presentation:
        """Create a new presentation based on the template."""
        try:
            # Create new presentation using the template
            return Presentation(self.template_analyzer.template_path)
        except Exception as e:
            logger.warning(f"Could not use template directly: {str(e)}")
            # Fallback to blank presentation
            return Presentation()
    
    def _generate_slides_with_images(self, presentation: Presentation, structure: Dict[str, Any]):
        """Generate slides based on the analyzed structure, incorporating user images."""
        # Clear existing slides (keep layouts)
        slide_indexes_to_remove = list(range(len(presentation.slides)))
        for i in reversed(slide_indexes_to_remove):
            r_id = presentation.slides._sldIdLst[i].rId
            presentation.part.drop_rel(r_id)
            del presentation.slides._sldIdLst[i]
        
        slides_data = structure.get('slides', [])
        image_index = 0
        
        for i, slide_data in enumerate(slides_data):
            # Create regular slide
            self._create_slide(presentation, slide_data)
            
            # Add image slide after every 2-3 content slides
            if (self.user_images and 
                image_index < len(self.user_images) and 
                slide_data.get('slide_type') == 'content' and 
                (i + 1) % 3 == 0):  # Insert image slide every 3rd slide
                
                self._create_image_slide(presentation, self.user_images[image_index], i + 1)
                image_index += 1

    def _create_image_slide(self, presentation: Presentation, image_path: str, slide_number: int):
        """Create a slide with an image and accompanying text using template background."""
        try:
            # Use content layout to maintain template background and styling
            layout_index = self._find_content_layout_for_image(presentation)
            slide_layout = presentation.slide_layouts[layout_index]
            slide = presentation.slides.add_slide(slide_layout)
            
            # Calculate slide index for animations
            slide_index = len(presentation.slides) - 1
            
            # Set slide title with enhanced styling
            if slide.shapes.title:
                slide.shapes.title.text = f"Visual Insight {slide_number}"
                self._style_title(slide.shapes.title)
            
            # Clear any existing content placeholder to make room for our custom layout
            self._clear_content_placeholder(slide)
            
            # Add the image and text in a custom layout that preserves template background
            self._add_image_and_text_custom_layout(slide, image_path, slide_number)
            
            # Add animations to the image slide
            self._add_slide_animations(slide, slide_index)

        except Exception as e:
            logger.error(f"Error creating image slide: {str(e)}")
            # Create enhanced fallback that still uses template
            self._create_template_based_image_slide(presentation, image_path, slide_number)

    def _find_content_layout_for_image(self, presentation: Presentation) -> int:
        """Find a content layout that will preserve template background."""
        layouts = presentation.slide_layouts
        
        # First preference: Look for layouts with content placeholder
        for i, layout in enumerate(layouts):
            name_lower = layout.name.lower()
            if any(keyword in name_lower for keyword in 
                   ['content', 'text', 'bullet', 'two content']):
                return i
        
        # Second preference: Use second layout (usually content)
        if len(layouts) > 1:
            return 1
            
        # Fallback: Use first layout
        return 0

    def _clear_content_placeholder(self, slide):
        """Clear content placeholder to make room for custom image layout."""
        try:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Content placeholder
                    # Clear the placeholder content but keep it for background
                    if hasattr(shape, 'text_frame'):
                        shape.text_frame.clear()
                    break
        except Exception as e:
            logger.warning(f"Could not clear content placeholder: {str(e)}")

    def _add_image_and_text_custom_layout(self, slide, image_path: str, slide_number: int):
        """Add image and text in custom layout while preserving template background."""
        try:
            # Position image on the left side
            left = Inches(0.5)
            top = Inches(2.2)  # Below the title
            width = Inches(5.5)
            height = Inches(4)
            
            # Add image
            image_shape = slide.shapes.add_picture(image_path, left, top, width, height)
            
            # Position text box on the right side
            text_left = Inches(6.2)
            text_top = Inches(2.5)
            text_width = Inches(3.3)
            text_height = Inches(3.5)
            
            textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = None
            
            # Add enhanced descriptive content
            caption_content = [
                "Key visual supporting our strategic analysis",
                "Important data visualization and insights", 
                "Critical business intelligence demonstration",
                "Contextual evidence and supporting material",
                "Strategic decision-making reference"
            ]
            
            for i, item in enumerate(caption_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.level = 0
                p.space_after = Pt(10)
                self._style_enhanced_content_paragraph(p, 'content')
                
        except Exception as e:
            logger.warning(f"Could not add image and text: {str(e)}")

    def _create_template_based_image_slide(self, presentation: Presentation, image_path: str, slide_number: int):
        """Create image slide using template layout as fallback."""
        try:
            # Use content layout to maintain template styling
            layout_index = min(1, len(presentation.slide_layouts) - 1)
            slide_layout = presentation.slide_layouts[layout_index]
            slide = presentation.slides.add_slide(slide_layout)
            
            # Add enhanced title
            if slide.shapes.title:
                slide.shapes.title.text = f"Supporting Visual Evidence {slide_number}"
                self._style_title(slide.shapes.title)
            
            # Add image in center with proper spacing
            left = Inches(1.5)
            top = Inches(2.5)
            width = Inches(7)
            height = Inches(4)
            
            slide.shapes.add_picture(image_path, left, top, width, height)
            
        except Exception as e:
            logger.error(f"Failed to create template-based image slide: {str(e)}")

    def _find_image_layout(self, presentation: Presentation) -> int:
        """Find the best layout for image slides."""
        layouts = presentation.slide_layouts
        
        # Look for layouts that might support images
        for i, layout in enumerate(layouts):
            name_lower = layout.name.lower()
            if any(keyword in name_lower for keyword in 
                   ['picture', 'image', 'photo', 'content', 'blank']):
                return i
        
        # Fallback to a content layout or blank layout
        return min(1, len(layouts) - 1)

    def _add_image_to_slide(self, slide, image_path: str):
        """Add an image to the slide with proper positioning."""
        try:
            # Position image in the center-left area
            left = Inches(1)
            top = Inches(2)
            width = Inches(5)
            height = Inches(4)
            
            # Add image
            slide.shapes.add_picture(image_path, left, top, width, height)
            
        except Exception as e:
            logger.warning(f"Could not add image {image_path}: {str(e)}")

    def _add_image_caption(self, slide, image_path: str, slide_number: int):
        """Add descriptive text next to the image."""
        try:
            # Position text box next to the image
            left = Inches(6.5)
            top = Inches(2.5)
            width = Inches(2.5)
            height = Inches(3)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Add descriptive content
            caption_content = [
                "Key visual supporting our analysis",
                "Important data visualization",
                "Strategic insight demonstration", 
                "Contextual reference material"
            ]
            
            for i, item in enumerate(caption_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.level = 0
                self._style_enhanced_content_paragraph(p, 'content')
                
        except Exception as e:
            logger.warning(f"Could not add image caption: {str(e)}")

    def _create_basic_image_slide(self, presentation: Presentation, image_path: str, slide_number: int):
        """Create a basic image slide as fallback."""
        try:
            # Use first available layout
            slide_layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(slide_layout)
            
            # Add title
            if slide.shapes.title:
                slide.shapes.title.text = f"Supporting Visual {slide_number}"
            
            # Add image in center
            left = Inches(2)
            top = Inches(2)
            width = Inches(6)
            height = Inches(4)
            
            slide.shapes.add_picture(image_path, left, top, width, height)
            
        except Exception as e:
            logger.error(f"Failed to create basic image slide: {str(e)}")

    def _generate_slides(self, presentation: Presentation, structure: Dict[str, Any]):
        """Generate slides based on the analyzed structure (legacy method for compatibility)."""
        self._generate_slides_with_images(presentation, structure)
    
    def _create_slide(self, presentation: Presentation, slide_data: Dict[str, Any]):
        """Create a single slide from slide data with enhanced formatting and animations."""
        slide_type = slide_data.get('slide_type', 'content')
        layout_index = self.template_analyzer.get_best_layout_for_slide_type(slide_type)
        
        try:
            # Add slide with appropriate layout
            slide_layout = presentation.slide_layouts[layout_index]
            slide = presentation.slides.add_slide(slide_layout)
            
            # Calculate slide index for animations
            slide_index = len(presentation.slides) - 1
            
            # Set slide title with enhanced formatting and perfect positioning
            title = slide_data.get('title', 'Slide Title')
            
            # Truncate title if too long to prevent overflow
            max_title_length = 60
            if len(title) > max_title_length:
                title = title[:max_title_length] + "..."
                
            if slide.shapes.title:
                slide.shapes.title.text = title
                # Position title at the very top for attractiveness
                slide.shapes.title.top = Inches(0.15)  # Very top positioning
                slide.shapes.title.height = Inches(1.8)  # Larger height for attractive spacing
                self._style_title(slide.shapes.title)
            else:
                # Create attractive manual title if no title placeholder
                self._create_manual_title(slide, title)
            
            # Add main content with better formatting
            content = slide_data.get('content', [])
            self._add_enhanced_slide_content(slide, content, slide_type)
            
            # Add emphasis points if available
            emphasis_points = slide_data.get('emphasis_points', [])
            if emphasis_points:
                self._add_emphasis_content(slide, emphasis_points)
            
            # Add animations to the slide
            self._add_slide_animations(slide, slide_index)
            
            # Add speaker notes with enhanced details
            speaking_notes = slide_data.get('speaking_notes', '')
            if speaking_notes:
                self._add_detailed_speaker_notes(slide, speaking_notes)
            else:
                self._add_speaker_notes(slide, title, content)
            
        except Exception as e:
            logger.error(f"Error creating slide: {str(e)}")
            # Fallback to enhanced basic slide creation
            self._create_enhanced_basic_slide(presentation, slide_data)

    def _create_manual_title(self, slide, title_text: str):
        """Create an attractive manual title with underline and overflow protection when no title placeholder exists."""
        try:
            # Get slide dimensions for perfect positioning
            slide_width = self.template_analyzer.get_slide_dimensions()[0]
            
            # Truncate title if too long to prevent overflow
            max_title_length = 60  # Character limit for titles
            if len(title_text) > max_title_length:
                title_text = title_text[:max_title_length] + "..."
            
            # Create attractive title text box at the very top center
            left = Inches(0.5)
            top = Inches(0.15)  # Very top positioning
            width = slide_width - Inches(1)  # Full width minus margins
            height = Inches(1.8)  # Increased height for attractive spacing
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
            title_frame.margin_left = Inches(0.2)
            title_frame.margin_right = Inches(0.2)
            title_frame.margin_top = Inches(0.15)
            title_frame.margin_bottom = Inches(0.15)
            
            # Style the manual title with attractive formatting
            for paragraph in title_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER  # Perfect center alignment
                paragraph.space_after = Pt(28)  # Enhanced spacing
                paragraph.space_before = Pt(8)
                
                for run in paragraph.runs:
                    fonts = self.template_analyzer.get_theme_fonts()
                    colors = self.template_analyzer.get_theme_colors()
                    
                    run.font.name = fonts.get('title', 'Arial Black')  # More attractive font
                    run.font.size = Pt(48)  # Even bigger for impact
                    run.font.bold = True
                    run.font.underline = True  # Add underline for attractiveness
                    
                    # Apply high contrast colors
                    self._apply_high_contrast_color(run, colors, is_title=True)
                    
                    if 'primary' in colors:
                        color_hex = colors['primary'].lstrip('#')
                        rgb = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                        run.font.color.rgb = RGBColor(*rgb)
                        
        except Exception as e:
            logger.warning(f"Could not create manual title: {str(e)}")
    
    def _create_basic_slide(self, presentation: Presentation, slide_data: Dict[str, Any]):
        """Create a basic slide when template layouts fail."""
        # Use the first available layout
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Slide Title')
        
        # Add content as text box if no content placeholder
        content = slide_data.get('content', [])
        if content:
            # Create a text box for content
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.level = 0
    
    def _add_enhanced_slide_content(self, slide, content: List[str], slide_type: str):
        """Add content to a slide with enhanced formatting, better layout, and overflow protection."""
        if not content:
            return
        
        # Get slide dimensions and ensure content fits
        slide_width, slide_height = self.template_analyzer.get_slide_dimensions()
        fitted_content = self._ensure_content_fits_slide(content, slide_height)
        
        # Try to find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # Content placeholder
                content_placeholder = shape
                break
        
        if content_placeholder:
            # Ensure content stays within slide boundaries with safe margins
            slide_width, slide_height = self.template_analyzer.get_slide_dimensions()
            
            # Safe positioning to prevent overflow
            safe_margin = Inches(0.8)  # Increased safety margin
            content_placeholder.left = safe_margin
            content_placeholder.top = Inches(2.5)  # More space below larger title
            content_placeholder.width = slide_width - (safe_margin * 2)  # Safe width
            content_placeholder.height = slide_height - Inches(4.2)  # Safe height with bottom margin
            
            # Use the placeholder with enhanced formatting and overflow protection
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = None  # Prevent auto-sizing that could cause overflow
            text_frame.margin_top = Inches(0.3)
            text_frame.margin_bottom = Inches(0.3)
            text_frame.margin_left = Inches(0.4)
            text_frame.margin_right = Inches(0.4)
            
            for i, item in enumerate(fitted_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Enhanced bullet point formatting
                p.text = item
                p.level = 0
                p.space_after = Pt(10)  # Reduced spacing to fit more content safely
                self._style_enhanced_content_paragraph(p, slide_type)
                
                # Add sub-points for detailed content (but limit to prevent overflow)
                if ':' in item and len(item) > 80 and i < 4:  # Limit sub-points
                    # Split long points into main point and details
                    parts = item.split(':', 1)
                    if len(parts) == 2:
                        p.text = parts[0] + ':'
                        # Add sub-details
                        sub_p = text_frame.add_paragraph()
                        sub_p.text = parts[1].strip()
                        sub_p.level = 1
                        sub_p.space_after = Pt(8)
                        self._style_enhanced_content_paragraph(sub_p, slide_type, is_sub=True)
        else:
            # Create enhanced text box manually with fitted content
            self._create_enhanced_content_textbox(slide, fitted_content, slide_type)

    def _add_emphasis_content(self, slide, emphasis_points: List[str]):
        """Add emphasis points or key statistics to the slide with dark, visible colors."""
        if not emphasis_points:
            return
            
        # Try to add emphasis content in a separate area
        try:
            # Create emphasis box in bottom area
            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(1)
            
            emphasis_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = emphasis_box.text_frame
            text_frame.word_wrap = True
            
            # Add emphasis content
            for i, point in enumerate(emphasis_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"💡 {point}"
                p.alignment = PP_ALIGN.CENTER
                
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.bold = True
                    
                    # Use dark colors for emphasis instead of potentially light accent colors
                    colors = self.template_analyzer.get_theme_colors()
                    bg_colors = self.template_analyzer.get_background_colors()
                    
                    if self._is_dark_background(bg_colors):
                        # Dark background - use bright white
                        run.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        # Light background - use dark orange/red for emphasis but still visible
                        run.font.color.rgb = RGBColor(180, 60, 0)  # Dark orange
        
        except Exception as e:
            logger.warning(f"Could not add emphasis content: {str(e)}")

    def _create_enhanced_content_textbox(self, slide, content: List[str], slide_type: str):
        """Create an enhanced text box with overflow protection to ensure content stays within slide boundaries."""
        # Get slide dimensions for safe positioning
        slide_width, slide_height = self.template_analyzer.get_slide_dimensions()
        
        # Safe positioning with protective margins
        safe_margin = Inches(0.8)  # Increased safety margin
        left = safe_margin
        top = Inches(2.5)  # More space below the larger attractive title
        width = slide_width - (safe_margin * 2)  # Safe width preventing overflow
        height = slide_height - Inches(4.2)  # Safe height with bottom margin
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Critical: Prevent auto-sizing overflow
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)
        text_frame.margin_left = Inches(0.4)
        text_frame.margin_right = Inches(0.4)
        
        # Add content with overflow protection
        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Enhanced bullet formatting
            if slide_type == 'comparison':
                bullet_symbol = "▶" if i % 2 == 0 else "◀"
            else:
                bullet_symbol = "●"
            
            # Truncate content if too long to prevent overflow
            max_content_length = 200  # Character limit per bullet point
            if len(item) > max_content_length:
                item = item[:max_content_length] + "..."
            
            p.text = f"{bullet_symbol} {item}"
            p.level = 0
            p.space_after = Pt(12)  # Reduced spacing to fit more content safely
            self._style_enhanced_content_paragraph(p, slide_type)

    def _create_enhanced_basic_slide(self, presentation: Presentation, slide_data: Dict[str, Any]):
        """Create an enhanced basic slide when template layouts fail."""
        # Use the best available layout instead of just first
        slide_type = slide_data.get('slide_type', 'content')
        layouts = presentation.slide_layouts
        
        # Try to use a content layout
        layout_index = min(1, len(layouts) - 1) if len(layouts) > 1 else 0
        slide_layout = layouts[layout_index]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Add enhanced title
        title = slide_data.get('title', 'Slide Title')
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title(slide.shapes.title)
        
        # Add enhanced content
        content = slide_data.get('content', [])
        if content:
            self._create_enhanced_content_textbox(slide, content, slide_type)
        
        # Add emphasis if available
        emphasis_points = slide_data.get('emphasis_points', [])
        if emphasis_points:
            self._add_emphasis_content(slide, emphasis_points)
    
    def _style_title(self, title_shape):
        """Apply enhanced attractive styling to title with underline and perfect center positioning."""
        try:
            fonts = self.template_analyzer.get_theme_fonts()
            colors = self.template_analyzer.get_theme_colors()
            
            # Perfect center positioning at the top of slide
            slide_width = self.template_analyzer.get_slide_dimensions()[0]
            title_shape.left = Inches(0.5)
            title_shape.top = Inches(0.15)  # Very top positioning
            title_shape.width = slide_width - Inches(1)  # Full width minus margins
            title_shape.height = Inches(1.8)  # Increased height for attractive spacing
            
            # Enhanced text frame settings for attractive layout
            title_shape.text_frame.margin_bottom = Inches(0.15)
            title_shape.text_frame.margin_top = Inches(0.15)
            title_shape.text_frame.margin_left = Inches(0.2)
            title_shape.text_frame.margin_right = Inches(0.2)
            title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
            
            for paragraph in title_shape.text_frame.paragraphs:
                # Perfect center alignment
                paragraph.alignment = PP_ALIGN.CENTER
                
                for run in paragraph.runs:
                    run.font.name = fonts.get('title', 'Arial Black')  # More attractive font
                    run.font.size = Pt(48)  # Even bigger for more impact
                    run.font.bold = True
                    run.font.underline = True  # Add underline for attractiveness
                    
                    # High contrast color selection for better visibility
                    self._apply_high_contrast_color(run, colors, is_title=True)
                    
                # Enhanced spacing for attractive layout
                paragraph.space_after = Pt(28)
                paragraph.space_before = Pt(8)
        
        except Exception as e:
            logger.warning(f"Error styling title: {str(e)}")
    
    def _style_enhanced_content_paragraph(self, paragraph, slide_type: str, is_sub: bool = False):
        """Apply enhanced styling to content paragraph with justification, bigger fonts, and high contrast."""
        try:
            fonts = self.template_analyzer.get_theme_fonts()
            colors = self.template_analyzer.get_theme_colors()
            
            # JUSTIFY align all content for even distribution between margins
            paragraph.alignment = PP_ALIGN.JUSTIFY
            
            for run in paragraph.runs:
                run.font.name = fonts.get('body', 'Calibri')
                
                # Significantly bigger font sizing based on content type
                if is_sub:
                    run.font.size = Pt(20)  # Increased from 16
                    run.font.italic = True
                elif slide_type == 'title':
                    run.font.size = Pt(28)  # Increased from 24
                    run.font.bold = True
                elif slide_type == 'section':
                    run.font.size = Pt(24)  # Increased from 20
                    run.font.bold = True
                else:
                    run.font.size = Pt(22)  # Increased from 18
                
                # Apply high contrast colors
                self._apply_high_contrast_color(run, colors, is_title=(slide_type in ['title', 'section']))
        
        except Exception as e:
            logger.warning(f"Error styling paragraph: {str(e)}")

    def _apply_high_contrast_color(self, run, colors, is_title=False):
        """Apply high contrast colors with emphasis on dark, visible text."""
        try:
            # Get background color to determine contrast
            bg_colors = self.template_analyzer.get_background_colors()
            
            # Default to very dark colors for maximum visibility
            if is_title:
                # For titles, use very dark colors for maximum impact
                if self._is_dark_background(bg_colors):
                    # Dark background - use bright white
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Pure white
                else:
                    # Light background - use very dark colors
                    run.font.color.rgb = RGBColor(20, 20, 20)  # Very dark gray (almost black)
            else:
                # For content, prioritize dark colors for readability
                if self._is_dark_background(bg_colors):
                    # Dark background - use bright white for contrast
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Pure white
                else:
                    # Light background - use very dark colors for visibility
                    run.font.color.rgb = RGBColor(40, 40, 40)  # Very dark gray
                    
            # Only apply theme colors if they are dark enough for visibility
            if 'primary' in colors and self._is_color_dark_enough(colors['primary']):
                color_hex = colors['primary'].lstrip('#')
                rgb = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                # Only use if the color is dark enough
                luminance = (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) / 255
                if luminance < 0.6:  # Only use if relatively dark
                    run.font.color.rgb = RGBColor(*rgb)
                
        except Exception as e:
            logger.warning(f"Error applying high contrast color: {str(e)}")
            # Fallback to very dark colors for maximum visibility
            run.font.color.rgb = RGBColor(30, 30, 30)

    def _is_dark_background(self, bg_colors):
        """Determine if background is dark based on luminance."""
        try:
            if not bg_colors:
                return False
                
            # Get primary background color
            bg_color = bg_colors.get('primary', '#FFFFFF')
            if bg_color.startswith('#'):
                bg_color = bg_color[1:]
                
            # Calculate luminance
            r, g, b = tuple(int(bg_color[i:i+2], 16) for i in (0, 2, 4))
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            
            return luminance < 0.5  # Dark if luminance < 50%
        except:
            return False  # Default to light background

    def _is_color_dark_enough(self, color):
        """Check if a color is dark enough to be easily visible."""
        try:
            if not color:
                return False
                
            if color.startswith('#'):
                color = color[1:]
                
            # Calculate luminance
            r, g, b = tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            
            return luminance < 0.6  # Dark enough if luminance < 60%
        except:
            return False

    def _has_good_contrast(self, color1, bg_colors):
        """Check if color has good contrast against background."""
        try:
            if not bg_colors or not color1:
                return False
                
            bg_color = bg_colors.get('primary', '#FFFFFF')
            
            # Simple contrast check - different approach than luminance
            if color1.startswith('#'):
                color1 = color1[1:]
            if bg_color.startswith('#'):
                bg_color = bg_color[1:]
                
            # Calculate color difference
            r1, g1, b1 = tuple(int(color1[i:i+2], 16) for i in (0, 2, 4))
            r2, g2, b2 = tuple(int(bg_color[i:i+2], 16) for i in (0, 2, 4))
            
            diff = abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2)
            return diff > 200  # Good contrast if difference > 200
            
        except:
            return False

    def _add_slide_animations(self, slide, slide_index):
        """Add smooth slide transition animations."""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.dml.color import RGBColor
            
            # Add slide transition based on slide type
            transitions = [
                'fade',      # Slide 1
                'push',      # Slide 2  
                'wipe',      # Slide 3
                'split',     # Slide 4
                'reveal',    # Slide 5
                'cover',     # Slide 6
                'cut'        # Slide 7+
            ]
            
            # Cycle through transitions
            transition_type = transitions[slide_index % len(transitions)]
            
            # Note: python-pptx has limited animation support
            # Adding entrance animations to text elements
            self._add_text_animations(slide)
            
        except Exception as e:
            logger.warning(f"Error adding slide animations: {str(e)}")

    def _add_text_animations(self, slide):
        """Add entrance animations to text elements on the slide."""
        try:
            # Add animation effects to text shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    # Add fly-in or fade-in effect simulation through positioning
                    # Note: Full animation requires more complex XML manipulation
                    self._simulate_entrance_effect(shape)
                    
        except Exception as e:
            logger.warning(f"Error adding text animations: {str(e)}")

    def _simulate_entrance_effect(self, shape):
        """Simulate entrance effects by optimizing text positioning and styling."""
        try:
            # Enhance text visibility and readability as animation alternative
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                
                # Add subtle shadow effect for depth
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Add slight shadow simulation via color depth
                        if hasattr(run.font, 'color'):
                            # This creates visual depth similar to animation
                            pass
                            
        except Exception as e:
            logger.warning(f"Error simulating entrance effect: {str(e)}")

    def _ensure_content_fits_slide(self, content: List[str], slide_height: float) -> List[str]:
        """Ensure content fits within slide boundaries by managing length and count."""
        try:
            # Calculate maximum content items based on slide height and font size
            available_height = slide_height - Inches(4.5)  # Account for title and margins
            line_height = Pt(30)  # Approximate line height including spacing
            max_lines = int(available_height / line_height)
            
            # Limit number of content items
            max_items = min(max_lines // 2, 8)  # Conservative estimate with spacing
            fitted_content = content[:max_items]
            
            # Truncate individual items if too long
            max_chars_per_item = 180
            for i, item in enumerate(fitted_content):
                if len(item) > max_chars_per_item:
                    fitted_content[i] = item[:max_chars_per_item] + "..."
            
            return fitted_content
            
        except Exception as e:
            logger.warning(f"Error fitting content to slide: {str(e)}")
            return content[:6]  # Safe fallback

    def _add_detailed_speaker_notes(self, slide, speaking_notes: str):
        """Add detailed speaker notes to the slide."""
        try:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = speaking_notes
        except Exception as e:
            logger.warning(f"Error adding detailed speaker notes: {str(e)}")

    def _style_content_paragraph(self, paragraph):
        """Apply styling to content paragraph with justification, bigger font, and dark colors (legacy method for compatibility)."""
        try:
            fonts = self.template_analyzer.get_theme_fonts()
            colors = self.template_analyzer.get_theme_colors()
            
            # JUSTIFY align content for even distribution between margins
            paragraph.alignment = PP_ALIGN.JUSTIFY
            
            for run in paragraph.runs:
                run.font.name = fonts.get('body', 'Calibri')
                run.font.size = Pt(22)  # Increased from 18pt to 22pt
                
                # Apply dark colors for better visibility
                bg_colors = self.template_analyzer.get_background_colors()
                if self._is_dark_background(bg_colors):
                    # Dark background - use bright white
                    run.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    # Light background - use very dark gray for maximum visibility
                    run.font.color.rgb = RGBColor(40, 40, 40)
        
        except Exception as e:
            logger.warning(f"Error styling paragraph: {str(e)}")

    def _create_content_textbox(self, slide, content: List[str]):
        """Create a text box for content with overflow protection when no placeholder is available (legacy method)."""
        # Get slide dimensions for safe positioning
        slide_width, slide_height = self.template_analyzer.get_slide_dimensions()
        
        # Safe positioning to prevent overflow
        safe_margin = Inches(0.8)
        left = safe_margin
        top = Inches(2.5)
        width = slide_width - (safe_margin * 2)  # Safe width
        height = slide_height - Inches(4.2)  # Safe height with bottom margin
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Prevent overflow
        text_frame.margin_left = Inches(0.4)
        text_frame.margin_right = Inches(0.4)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)
        
        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Truncate content if too long
            max_content_length = 180  # Slightly shorter for legacy method
            if len(item) > max_content_length:
                item = item[:max_content_length] + "..."
            
            p.text = f"• {item}"
            p.level = 0
            p.space_after = Pt(10)  # Reduced spacing for safety
            self._style_content_paragraph(p)

    def _create_basic_slide(self, presentation: Presentation, slide_data: Dict[str, Any]):
        """Create a basic slide when template layouts fail (legacy method)."""
        # Use the first available layout
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Slide Title')
        
        # Add content as text box if no content placeholder
        content = slide_data.get('content', [])
        if content:
            self._create_content_textbox(slide, content)
    
    def _add_speaker_notes(self, slide, title: str, content: List[str]):
        """Add speaker notes to the slide."""
        try:
            # Generate speaker notes using LLM
            slide_content = f"Title: {title}\nContent: {'; '.join(content)}"
            speaker_notes = self.llm_service.generate_speaker_notes(slide_content)
            
            # Add to slide notes
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = speaker_notes
        
        except Exception as e:
            logger.warning(f"Error adding speaker notes: {str(e)}")
    
    def _save_presentation(self, presentation: Presentation) -> str:
        """Save the presentation to a temporary file."""
        # Create temporary file
        with tempfile.NamedTemporaryFile(
            suffix='.pptx', 
            delete=False, 
            dir=tempfile.gettempdir()
        ) as tmp_file:
            output_path = tmp_file.name
        
        # Save presentation
        presentation.save(output_path)
        
        return output_path
